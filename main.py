import csv
import os
import random
import time

from openpyxl.utils import get_column_letter
from xlrd import open_workbook
from xlutils.copy import copy
import win32com
from docx import Document
from docx.shared import Cm, Inches
from openpyxl import Workbook
from win32com.client import Dispatch
from pptx import Presentation

from mkdir import dofile
from setting import *

print(time.ctime())

basepath = os.path.dirname(os.path.abspath(__file__))
filepath = os.path.join(basepath, "basefiles")
# 创建文件目录
dofile()

# ------------------------------------------------------------------------------
# 生成xls
XLS_FILE_EXT = ".xls"
namelist_XLSX = [str(temp) + XLS_FILE_EXT for temp in range(1, XLS_NUM + 1)]
# rb = open_workbook(r"E:\PycharmProjects\test\fileCreat\1.xls")
#
# # 通过sheet_by_index()获取的sheet没有write()方法
# rs = rb.sheet_by_index(0)
# wb = copy(rb)
# ws = wb.get_sheet(0)
xls_csv_list = []
for i in range(XLS_NUM):
    rb = open_workbook(os.path.join(filepath, "1.xls"))
    # 通过sheet_by_index()获取的sheet没有write()方法
    rs = rb.sheet_by_index(0)
    wb = copy(rb)
    ws = wb.get_sheet(0)
    ws.write(0, 0, str(i))
    # 从第2行开始，写入9行10列数据，值为对应的列序号A、B、C、D...
    for row in range(2, XLS_MAXROW):
        for col in range(1, XLS_MAXCOLUMN):
            ws.write(row, col, get_column_letter(col))

    # 可以使用append插入一行数据
    ws.write(XLS_MAXROW, 1, str(int(round(time.time() * 1000))) + str(
        random.random()))
    wb.save(XLS_SAVE_PATH + '/' + XLS_NAME_SET + namelist_XLSX[i])
    xls_csv_list.append([XLS_SAVE_PATH + '/' + XLS_NAME_SET + namelist_XLSX[i]])
with open(XLS_SAVE_PATH + '/' + XLS_NAME_SET[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in xls_csv_list:
        csv_write.writerow(temp)
f.close()
print("--------------------------------------------------")
print(">> excel for xls 任务执行完成...")
print(time.ctime())
# ------------------------------------------------------------------------------
# 生成xlsx
XLSX_FILE_EXT = ".xlsx"
namelist_XLSX = [str(temp) + XLSX_FILE_EXT for temp in range(1, XLSX_NUM + 1)]
# 在内存中创建一个workbook对象，而且会至少创建一个 worksheet
wb = Workbook()

# 获取当前活跃的worksheet,默认就是第一个worksheet
ws = wb.active

# 设置单元格的值，A1等于6(测试可知openpyxl的行和列编号从1开始计算)，B1等于7
xlsx_csv_list = []
for i in range(XLSX_NUM):
    ws.cell(row=1, column=1).value = str(i)
    # ws.cell(row=1, column=2).value = 7
    # 从第2行开始，写入9行10列数据，值为对应的列序号A、B、C、D...
    for row in range(2, XLSX_MAXROW):
        for col in range(1, XLSX_MAXCOLUMN):
            ws.cell(row=row, column=col).value = get_column_letter(col)

    # 可以使用append插入一行数据
    ws.cell(row=XLSX_MAXROW, column=1).value = str(int(round(time.time() * 1000))) + str(
        random.random())
    wb.save(XLSX_SAVE_PATH + '/' + XLSX_NAME_SET + namelist_XLSX[i])
    xlsx_csv_list.append([XLSX_SAVE_PATH + '/' + XLSX_NAME_SET + namelist_XLSX[i]])
with open(XLSX_SAVE_PATH + '/' + XLSX_NAME_SET[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in xlsx_csv_list:
        csv_write.writerow(temp)
f.close()
print(">> excel for xlsx 任务执行完成...")
print(time.ctime())
# ------------------------------------------------------------------------------
# 生成doc
DOC_FILE_EXT = ".doc"
namelist_DOC = [str(temp) + DOC_FILE_EXT for temp in range(1, DOC_NUM + 1)]

# textdoc
# 填充content
content = ['在'
           '一个文具店里，有一个美丽的文具盒，'
           '铅笔、橡皮、转笔刀。'
           '呆在货架上时间久了，'
           '也就成为好朋友了。他'
           '们共同渴望着自己的主'
           '人来把它们买回家。终'
           '于，小姑娘小梅把它们买了回小梅上一年级，因为作业上需要它们，所以催着妈'
           '妈就把它们'
           '买了回来。铅笔、橡皮、转笔刀三个好朋友在小梅家'
           '住了一段时间，都'
           '慢慢骄傲起来，越来越瞧'
           '不起对方。铅笔自满'
           '地说：“哼，我的功劳最'
           '大，你们两个都是废物。”橡皮说：“才怪，你写错了字，还不是我来'
           '帮你擦，要不然，主人的作业'
           '脏死了了。”转笔'
           '刀不服气了：“难道我'
           '就没有功劳了吗?铅笔，你这'
           '个忘恩负义的家伙。”转笔刀流下了伤心的泪水：“你'
           '断的时候还不是我帮'
           '你转，要不然主人'
           '早就不要你了。”文具'
           '盒劝它们说：“你们不要吵了。你们都是小主人最需要的东西，你们争吵是'
           '没用的，只会给主人'
           '带来不方便、不愉'
           '快。”铅笔、转笔刀，橡皮'
           '惭愧地低下了头。从此，它们和睦相处，从不争'
           '吵帮着主人好好学习。']
DOC_NUM1_CSV_list = []
for i in range(DOC_NUM):
    myDocument = Document()
    for j in range(DOC_PAGE_NUM1):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH1 + "/" + DOC_NAME_SET1 + namelist_DOC[i])
    DOC_NUM1_CSV_list.append([DOC_SAVE_PATH1 + "/" + DOC_NAME_SET1 + namelist_DOC[i]])
with open(DOC_SAVE_PATH1 + "/" + DOC_NAME_SET1[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in DOC_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()
DOC_NUM2_CSV_list = []
for i in range(DOC_NUM):
    myDocument = Document()
    for j in range(DOC_PAGE_NUM2):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH2 + "/" + DOC_NAME_SET2 + namelist_DOC[i])
    DOC_NUM2_CSV_list.append([DOC_SAVE_PATH2 + "/" + DOC_NAME_SET2 + namelist_DOC[i]])
with open(DOC_SAVE_PATH2 + "/" + DOC_NAME_SET2[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in DOC_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()
DOC_NUM3_CSV_list = []
for i in range(DOC_NUM):
    myDocument = Document()
    for j in range(DOC_PAGE_NUM3):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH3 + "/" + DOC_NAME_SET3 + namelist_DOC[i])
    DOC_NUM3_CSV_list.append([DOC_SAVE_PATH3 + "/" + DOC_NAME_SET3 + namelist_DOC[i]])
with open(DOC_SAVE_PATH3 + "/" + DOC_NAME_SET3[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in DOC_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()
DOC_NUM4_CSV_list = []
for i in range(DOC_NUM):
    myDocument = Document()
    for j in range(DOC_PAGE_NUM4):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH4 + "/" + DOC_NAME_SET4 + namelist_DOC[i])
    DOC_NUM4_CSV_list.append([DOC_SAVE_PATH4 + "/" + DOC_NAME_SET4 + namelist_DOC[i]])
with open(DOC_SAVE_PATH4 + "/" + DOC_NAME_SET4[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in DOC_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

# picdoc
PICDOC_NUM1_CSV_list = []
namelist_DOC = [str(temp) + DOC_FILE_EXT for temp in range(1, DOC_NUM + 1)]
for a in range(DOC_NUM):
    myDocument = Document()
    for i in range(DOC_PAGE_NUM1):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH11 + "/" + DOC_NAME_SET11 + namelist_DOC[a])
    PICDOC_NUM1_CSV_list.append([DOC_SAVE_PATH11 + "/" + DOC_NAME_SET11 + namelist_DOC[a]])
with open(DOC_SAVE_PATH11 + "/" + DOC_NAME_SET11[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOC_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOC_NUM2_CSV_list = []
for a in range(DOC_NUM):
    myDocument = Document()
    for i in range(DOC_PAGE_NUM2):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH12 + "/" + DOC_NAME_SET12 + namelist_DOC[a])
    PICDOC_NUM2_CSV_list.append([DOC_SAVE_PATH12 + "/" + DOC_NAME_SET12 + namelist_DOC[a]])
with open(DOC_SAVE_PATH12 + "/" + DOC_NAME_SET12[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOC_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOC_NUM3_CSV_list = []
for a in range(DOC_NUM):
    myDocument = Document()
    for i in range(DOC_PAGE_NUM3):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH13 + "/" + DOC_NAME_SET13 + namelist_DOC[a])
    PICDOC_NUM3_CSV_list.append([DOC_SAVE_PATH13 + "/" + DOC_NAME_SET13 + namelist_DOC[a]])
with open(DOC_SAVE_PATH13 + "/" + DOC_NAME_SET13[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOC_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOC_NUM4_CSV_list = []
for a in range(DOC_NUM):
    myDocument = Document()
    for i in range(DOC_PAGE_NUM4):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOC_SAVE_PATH14 + "/" + DOC_NAME_SET14 + namelist_DOC[a])
    PICDOC_NUM4_CSV_list.append([DOC_SAVE_PATH14 + "/" + DOC_NAME_SET14 + namelist_DOC[a]])
with open(DOC_SAVE_PATH14 + "/" + DOC_NAME_SET14[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOC_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

print(">> word for doc 任务执行完成...")
print(time.ctime())
# ------------------------------------------------------------------------------
# 生成docx
DOCX_FILE_EXT = ".docx"
namelist_DOCX = [str(temp) + DOCX_FILE_EXT for temp in range(1, DOCX_NUM + 1)]

# textdoc
# 填充content
content = ['在'
           '一个文具店里，有一个美丽的文具盒，'
           '铅笔、橡皮、转笔刀。'
           '呆在货架上时间久了，'
           '也就成为好朋友了。他'
           '们共同渴望着自己的主'
           '人来把它们买回家。终'
           '于，小姑娘小梅把它们买了回小梅上一年级，因为作业上需要它们，所以催着妈'
           '妈就把它们'
           '买了回来。铅笔、橡皮、转笔刀三个好朋友在小梅家'
           '住了一段时间，都'
           '慢慢骄傲起来，越来越瞧'
           '不起对方。铅笔自满'
           '地说：“哼，我的功劳最'
           '大，你们两个都是废物。”橡皮说：“才怪，你写错了字，还不是我来'
           '帮你擦，要不然，主人的作业'
           '脏死了了。”转笔'
           '刀不服气了：“难道我'
           '就没有功劳了吗?铅笔，你这'
           '个忘恩负义的家伙。”转笔刀流下了伤心的泪水：“你'
           '断的时候还不是我帮'
           '你转，要不然主人'
           '早就不要你了。”文具'
           '盒劝它们说：“你们不要吵了。你们都是小主人最需要的东西，你们争吵是'
           '没用的，只会给主人'
           '带来不方便、不愉'
           '快。”铅笔、转笔刀，橡皮'
           '惭愧地低下了头。从此，它们和睦相处，从不争'
           '吵帮着主人好好学习。']
TEXTDOCX_NUM1_CSV_list = []
for i in range(DOCX_NUM):
    myDocument = Document()
    for j in range(DOCX_PAGE_NUM1):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH1 + "/" + DOCX_NAME_SET1 + namelist_DOCX[i])
    TEXTDOCX_NUM1_CSV_list.append([DOCX_SAVE_PATH1 + "/" + DOCX_NAME_SET1 + namelist_DOCX[i]])
with open(DOCX_SAVE_PATH1 + "/" + DOCX_NAME_SET1[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTDOCX_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTDOCX_NUM2_CSV_list = []
for i in range(DOCX_NUM):
    myDocument = Document()
    for j in range(DOCX_PAGE_NUM2):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH2 + "/" + DOCX_NAME_SET2 + namelist_DOCX[i])
    TEXTDOCX_NUM2_CSV_list.append([DOCX_SAVE_PATH2 + "/" + DOCX_NAME_SET2 + namelist_DOCX[i]])
with open(DOCX_SAVE_PATH2 + "/" + DOCX_NAME_SET2[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTDOCX_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTDOCX_NUM3_CSV_list = []
for i in range(DOCX_NUM):
    myDocument = Document()
    for j in range(DOCX_PAGE_NUM3):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH3 + "/" + DOCX_NAME_SET3 + namelist_DOCX[i])
    TEXTDOCX_NUM3_CSV_list.append([DOCX_SAVE_PATH3 + "/" + DOCX_NAME_SET3 + namelist_DOCX[i]])
with open(DOCX_SAVE_PATH3 + "/" + DOCX_NAME_SET3[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTDOCX_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTDOCX_NUM4_CSV_list = []
for i in range(DOCX_NUM):
    myDocument = Document()
    for j in range(DOCX_PAGE_NUM4):
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(content[0])
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
        myDocument.add_paragraph(str(str(i) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH4 + "/" + DOCX_NAME_SET4 + namelist_DOCX[i])
    TEXTDOCX_NUM4_CSV_list.append([DOCX_SAVE_PATH4 + "/" + DOCX_NAME_SET4 + namelist_DOCX[i]])
with open(DOCX_SAVE_PATH4 + "/" + DOCX_NAME_SET4[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTDOCX_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

# picdoc
PICDOCX_NUM1_CSV_list = []
namelist_DOCX = [str(temp) + DOCX_FILE_EXT for temp in range(1, DOCX_NUM + 1)]
for a in range(DOCX_NUM):
    myDocument = Document()
    for i in range(DOCX_PAGE_NUM1):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH11 + "/" + DOCX_NAME_SET11 + namelist_DOCX[a])
    PICDOCX_NUM1_CSV_list.append([DOCX_SAVE_PATH11 + "/" + DOCX_NAME_SET11 + namelist_DOCX[a]])
with open(DOCX_SAVE_PATH11 + "/" + DOCX_NAME_SET11[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOCX_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOCX_NUM2_CSV_list = []
for a in range(DOCX_NUM):
    myDocument = Document()
    for i in range(DOCX_PAGE_NUM2):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH12 + "/" + DOCX_NAME_SET12 + namelist_DOCX[a])
    PICDOCX_NUM2_CSV_list.append([DOCX_SAVE_PATH12 + "/" + DOCX_NAME_SET12 + namelist_DOCX[a]])
with open(DOCX_SAVE_PATH12 + "/" + DOCX_NAME_SET12[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOCX_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOCX_NUM3_CSV_list = []
for a in range(DOCX_NUM):
    myDocument = Document()
    for i in range(DOCX_PAGE_NUM3):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH13 + "/" + DOCX_NAME_SET13 + namelist_DOCX[a])
    PICDOCX_NUM3_CSV_list.append([DOCX_SAVE_PATH13 + "/" + DOCX_NAME_SET13 + namelist_DOCX[a]])
with open(DOCX_SAVE_PATH13 + "/" + DOCX_NAME_SET13[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOCX_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

PICDOCX_NUM4_CSV_list = []
for a in range(DOCX_NUM):
    myDocument = Document()
    for i in range(DOCX_PAGE_NUM4):
        myDocument.add_picture(os.path.join(filepath, "12333.jpg"), width=Inches(6),
                               height=Inches(8))
    myDocument.add_paragraph(str(str(a) + str(int(round(time.time() * 1000)))))
    myDocument.save(DOCX_SAVE_PATH14 + "/" + DOCX_NAME_SET14 + namelist_DOCX[a])
    PICDOCX_NUM4_CSV_list.append([DOCX_SAVE_PATH14 + "/" + DOCX_NAME_SET14 + namelist_DOCX[a]])
with open(DOCX_SAVE_PATH14 + "/" + DOCX_NAME_SET14[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICDOCX_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

print(">> word for docx 任务执行完成...")
print(time.ctime())
# ------------------------------------------------------------------------------
# 生成ppt
PPT_FILE_EXT = ".ppt"
namelist_PPT = [str(temp) + PPT_FILE_EXT for temp in range(1, PPT_NUM + 1)]
# textpptx
TEXTPPTX_NUM1_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM1 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPT_SAVE_PATH1 + '/' + PPT_NAME_SET1 + namelist_PPT[i])
    TEXTPPTX_NUM1_CSV_list.append([PPT_SAVE_PATH1 + '/' + PPT_NAME_SET1 + namelist_PPT[i]])
with open(PPT_SAVE_PATH1 + '/' + PPT_NAME_SET1[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTX_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTX_NUM2_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM2 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPT_SAVE_PATH2 + '/' + PPT_NAME_SET2 + namelist_PPT[i])
    TEXTPPTX_NUM2_CSV_list.append([PPT_SAVE_PATH2 + '/' + PPT_NAME_SET2 + namelist_PPT[i]])
with open(PPT_SAVE_PATH2 + '/' + PPT_NAME_SET2[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTX_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTX_NUM3_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM3 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPT_SAVE_PATH3 + '/' + PPT_NAME_SET3 + namelist_PPT[i])
    TEXTPPTX_NUM3_CSV_list.append([PPT_SAVE_PATH3 + '/' + PPT_NAME_SET3 + namelist_PPT[i]])
with open(PPT_SAVE_PATH3 + '/' + PPT_NAME_SET3[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTX_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTX_NUM4_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM4 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPT_SAVE_PATH4 + '/' + PPT_NAME_SET4 + namelist_PPT[i])
    TEXTPPTX_NUM4_CSV_list.append([PPT_SAVE_PATH4 + '/' + PPT_NAME_SET4 + namelist_PPT[i]])
with open(PPT_SAVE_PATH4 + '/' + PPT_NAME_SET4[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTX_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

# picppt
PICPPT_NUM1_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM1 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPT_SAVE_PATH11 + '/' + PPT_NAME_SET11 + namelist_PPT[i])
    PICPPT_NUM1_CSV_list.append([PPT_SAVE_PATH11 + '/' + PPT_NAME_SET11 + namelist_PPT[i]])
with open(PPT_SAVE_PATH11 + '/' + PPT_NAME_SET11[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPT_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPT_NUM2_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM2 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPT_SAVE_PATH12 + '/' + PPT_NAME_SET12 + namelist_PPT[i])
    PICPPT_NUM2_CSV_list.append([PPT_SAVE_PATH12 + '/' + PPT_NAME_SET12 + namelist_PPT[i]])
with open(PPT_SAVE_PATH12 + '/' + PPT_NAME_SET12[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPT_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPT_NUM3_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM3 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPT_SAVE_PATH13 + '/' + PPT_NAME_SET13 + namelist_PPT[i])
    PICPPT_NUM3_CSV_list.append([PPT_SAVE_PATH13 + '/' + PPT_NAME_SET13 + namelist_PPT[i]])
with open(PPT_SAVE_PATH13 + '/' + PPT_NAME_SET13[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPT_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPT_NUM4_CSV_list = []
for i in range(PPT_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPT_PAGE_NUM4 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPT_SAVE_PATH14 + '/' + PPT_NAME_SET14 + namelist_PPT[i])
    PICPPT_NUM4_CSV_list.append([PPT_SAVE_PATH14 + '/' + PPT_NAME_SET14 + namelist_PPT[i]])
with open(PPT_SAVE_PATH14 + '/' + PPT_NAME_SET14[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPT_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

print(">> ppt for ppt 任务执行完成...")
print(time.ctime())
# ------------------------------------------------------------------------------
# 生成pptx
PPTX_FILE_EXT = ".pptx"
namelist_PPTX = [str(temp) + PPTX_FILE_EXT for temp in range(1, PPTX_NUM + 1)]
# textpptx
TEXTPPTXX_NUM1_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM1 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPTX_SAVE_PATH1 + '/' + PPTX_NAME_SET1 + namelist_PPTX[i])
    TEXTPPTXX_NUM1_CSV_list.append([PPTX_SAVE_PATH1 + '/' + PPTX_NAME_SET1 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH1 + '/' + PPTX_NAME_SET1[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTXX_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTXX_NUM2_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM2 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPTX_SAVE_PATH2 + '/' + PPTX_NAME_SET2 + namelist_PPTX[i])
    TEXTPPTXX_NUM2_CSV_list.append([PPTX_SAVE_PATH2 + '/' + PPTX_NAME_SET2 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH2 + '/' + PPTX_NAME_SET2[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTXX_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTXX_NUM3_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM3 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPTX_SAVE_PATH3 + '/' + PPTX_NAME_SET3 + namelist_PPTX[i])
    TEXTPPTXX_NUM3_CSV_list.append([PPTX_SAVE_PATH3 + '/' + PPTX_NAME_SET3 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH3 + '/' + PPTX_NAME_SET3[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTXX_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

TEXTPPTXX_NUM4_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM4 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(temp) + " Hello, World!" + str(int(round(time.time() * 1000)))
        subtitle.text = "测试组的老可爱们！"
    prs.save(PPTX_SAVE_PATH4 + '/' + PPTX_NAME_SET4 + namelist_PPTX[i])
    TEXTPPTXX_NUM4_CSV_list.append([PPTX_SAVE_PATH4 + '/' + PPTX_NAME_SET4 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH4 + '/' + PPTX_NAME_SET4[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in TEXTPPTXX_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()

# picpptx
PICPPTX_NUM1_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM1 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPTX_SAVE_PATH11 + '/' + PPTX_NAME_SET11 + namelist_PPTX[i])
    PICPPTX_NUM1_CSV_list.append([PPTX_SAVE_PATH11 + '/' + PPTX_NAME_SET11 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH11 + '/' + PPTX_NAME_SET11[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPTX_NUM1_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPTX_NUM2_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM2 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPTX_SAVE_PATH12 + '/' + PPTX_NAME_SET12 + namelist_PPTX[i])
    PICPPTX_NUM2_CSV_list.append([PPTX_SAVE_PATH12 + '/' + PPTX_NAME_SET12 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH12 + '/' + PPTX_NAME_SET12[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPTX_NUM2_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPTX_NUM3_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM3 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPTX_SAVE_PATH13 + '/' + PPTX_NAME_SET13 + namelist_PPTX[i])
    PICPPTX_NUM3_CSV_list.append([PPTX_SAVE_PATH13 + '/' + PPTX_NAME_SET13 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH13 + '/' + PPTX_NAME_SET13[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPTX_NUM3_CSV_list:
        csv_write.writerow(temp)
f.close()

PICPPTX_NUM4_CSV_list = []
for i in range(PPTX_NUM):
    prs = Presentation(os.path.join(filepath, "123.pptx"))
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(i) + " Hello, World!" + str(int(round(time.time() * 1000)))
    subtitle.text = "测试组的老可爱们！"
    for temp in range(PPTX_PAGE_NUM4 - 1):
        # PowerPoint()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        # shapes = slide.shapes
        slide.shapes.add_picture(os.path.join(filepath, "12333.jpg"), left=0, top=0)
        # slide.shapes.add_picture(PPTX_DEMO_PIC, left=0, top=0)
    prs.save(PPTX_SAVE_PATH14 + '/' + PPTX_NAME_SET14 + namelist_PPTX[i])
    PICPPTX_NUM4_CSV_list.append([PPTX_SAVE_PATH14 + '/' + PPTX_NAME_SET14 + namelist_PPTX[i]])
with open(PPTX_SAVE_PATH14 + '/' + PPTX_NAME_SET14[:-1] + '.csv', 'w', newline='') as f:
    csv_write = csv.writer(f)
    for temp in PICPPTX_NUM4_CSV_list:
        csv_write.writerow(temp)
f.close()
print(">> ppt for pptx 任务执行完成...")
print(time.ctime())
